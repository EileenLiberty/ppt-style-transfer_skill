# -*- coding: utf-8 -*-
"""v3: template TOC + readable fonts + fitted HD images on 16:9."""
from __future__ import annotations

import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path
from xml.sax.saxutils import escape

from pptx import Presentation

# Reuse geometry extraction from v2 builder
import importlib.util

_spec = importlib.util.spec_from_file_location(
    "image_geometry", Path(__file__).with_name("image_geometry.py")
)
_ig = importlib.util.module_from_spec(_spec)
assert _spec.loader
_spec.loader.exec_module(_ig)
extract_image_placements = _ig.extract_image_placements

ROOT = Path(r"g:\ppt generate")
TPL = ROOT / "模板-培训课件-20260703.pptx"
PPTX_SCRIPTS = ROOT / ".claude" / "skills" / "pptx" / "scripts"
# Filled by configure_paths()
SRC: Path
OUT: Path
WORK: Path
PNG_DIR: Path
SRC_UNPACKED: Path

FOOTER_RE = re.compile(r"^(中国民航大学\s*)?(CAUC)?$", re.I)
MASTER_JUNK = ("单击此处编辑",)

# Layout6 content band on 16:9 (below title chrome) — fill more of the page
CONTENT_BOX = {"x": 420000, "y": 1000000, "cx": 11350000, "cy": 5550000}

# Font sizes (hundredths of a point): 美观紧凑，忌过大与满页留白
SZ_COVER_TITLE = 3000  # 30pt
SZ_COVER_SUB = 1600  # 16pt
SZ_TITLE = 2200  # 22pt
SZ_BODY = 1500  # 15pt
SZ_BODY_IMG = 1400  # 14pt
SZ_TOC = 1500  # 15pt
SZ_SECTION = 2600  # 26pt


def xml_text(s: str) -> str:
    return escape(s, {"'": "&apos;", '"': "&quot;"})


def para(
    text: str,
    level: int = 0,
    bold: bool = False,
    sz: int = SZ_BODY,
    tight: bool = False,
) -> str:
    b = ' b="1"' if bold else ""
    # tighter paragraph spacing to reduce empty look
    spc = '<a:spcBef><a:spcPts val="60"/></a:spcBef><a:spcAft><a:spcPts val="60"/></a:spcAft>'
    if tight:
        spc = '<a:spcBef><a:spcPts val="40"/></a:spcBef><a:spcAft><a:spcPts val="40"/></a:spcAft>'
    return (
        f'<a:p><a:pPr lvl="{level}">{spc}</a:pPr>'
        f'<a:r><a:rPr lang="zh-CN" altLang="en-US" sz="{sz}"{b}/>'
        f"<a:t>{xml_text(text)}</a:t></a:r></a:p>"
    )


def empty_para() -> str:
    return '<a:p><a:endParaRPr lang="zh-CN" altLang="en-US"/></a:p>'


def make_pic(shape_id: int, name: str, rId: str, x: int, y: int, cx: int, cy: int) -> str:
    return f'''<p:pic>
  <p:nvPicPr>
    <p:cNvPr id="{shape_id}" name="{xml_text(name)}"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>
    <p:nvPr/>
  </p:nvPicPr>
  <p:blipFill>
    <a:blip r:embed="{rId}"/>
    <a:stretch><a:fillRect/></a:stretch>
  </p:blipFill>
  <p:spPr>
    <a:xfrm>
      <a:off x="{x}" y="{y}"/>
      <a:ext cx="{cx}" cy="{cy}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>'''


def ph_sp(shape_id: int, name: str, ph_attr: str, paragraphs_xml: str) -> str:
    return f'''<p:sp>
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{xml_text(name)}"/>
    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph {ph_attr}/></p:nvPr>
  </p:nvSpPr>
  <p:spPr/>
  <p:txBody>
    <a:bodyPr/><a:lstStyle/>
    {paragraphs_xml}
  </p:txBody>
</p:sp>'''


def title_box(shape_id: int, text: str, x=500000, y=200000, cx=11200000, cy=700000) -> str:
    """Free title for blank Layout9."""
    return f'''<p:sp>
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="标题"/>
    <p:cNvSpPr txBox="1"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{x}" y="{y}"/>
      <a:ext cx="{cx}" cy="{cy}"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"/><a:lstStyle/>
    {para(text, bold=True, sz=SZ_TITLE)}
  </p:txBody>
</p:sp>'''


def wrap_slide(parts: list[str]) -> str:
    tree = "\n".join(parts)
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
 xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
 xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/><a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      {tree}
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>
'''


def fit_placements(placements: list[dict]) -> list[dict]:
    """Keep relative layout; fit bounding box into 16:9 content area."""
    if not placements:
        return []
    minx = min(p["x"] for p in placements)
    miny = min(p["y"] for p in placements)
    maxx = max(p["x"] + p["cx"] for p in placements)
    maxy = max(p["y"] + p["cy"] for p in placements)
    bw, bh = max(maxx - minx, 1), max(maxy - miny, 1)
    scale = min(CONTENT_BOX["cx"] / bw, CONTENT_BOX["cy"] / bh) * 0.96
    ox = CONTENT_BOX["x"] + (CONTENT_BOX["cx"] - bw * scale) / 2
    oy = CONTENT_BOX["y"] + (CONTENT_BOX["cy"] - bh * scale) / 2
    out = []
    for p in placements:
        out.append(
            {
                **p,
                "x": int(ox + (p["x"] - minx) * scale),
                "y": int(oy + (p["y"] - miny) * scale),
                "cx": max(int(p["cx"] * scale), 10000),
                "cy": max(int(p["cy"] * scale), 10000),
            }
        )
    return out


def resolve_media_file(media_name: str, media_dir: Path, renamed: dict) -> str | None:
    """Prefer HD PNG for WMF; else src_* copy."""
    stem = Path(media_name).stem
    ext = Path(media_name).suffix.lower()
    if ext == ".wmf":
        png = PNG_DIR / f"{stem}.png"
        if png.exists() and png.stat().st_size > 500:
            dest = f"src_{stem}.png"
            if not (media_dir / dest).exists():
                shutil.copy2(png, media_dir / dest)
            return dest
    return renamed.get(media_name)


def extract_texts(prs: Presentation, idx: int, keep_org: bool = False) -> list[str]:
    slide = prs.slides[idx]
    items = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.replace("\r", "").strip()
        if not t or any(j in t for j in MASTER_JUNK):
            continue
        compact = re.sub(r"\s+", "", t)
        # Footer chrome on content pages; keep org name on cover
        if not keep_org and compact in ("中国民航大学CAUC", "中国民航大学", "CAUC"):
            continue
        if keep_org and compact == "CAUC":
            continue
        try:
            y = int(shape.top or 0)
            x = int(shape.left or 0)
        except Exception:
            y, x = 0, 0
        items.append((y, x, t))
    items.sort()
    lines = []
    for _, _, t in items:
        for line in t.split("\n"):
            line = line.strip()
            if line:
                lines.append(line)
    return lines


# Template TOC row geometry (from slideLayout4.xml) — generic for any chapter
TOC_PH_IDX = {
    3: ([1, 13, 14], "slideLayout2.xml"),
    4: ([1, 13, 14, 15], "slideLayout3.xml"),
    5: ([1, 13, 14, 15, 16], "slideLayout4.xml"),
}
# Layout4 标题占位符几何（幻灯片绝对坐标）
TOC_TITLE_X, TOC_TITLE_CX, TOC_TITLE_CY = 6936105, 4385310, 556895
TOC_ROW_YS = [1511935, 2338070, 3164205, 3990340, 4816475]  # Layout4 ph y
TOC_ROW_GAP = TOC_ROW_YS[1] - TOC_ROW_YS[0]  # 826135
# 序号框 / 标题描边框：由版式组合变换映射到幻灯片坐标（勿用子坐标系数值）
TOC_NUM_X, TOC_NUM_CX, TOC_NUM_CY = 5725160, 1016001, 523148
TOC_NUM_Y_OFF = 10252  # 相对标题 ph 的 y 偏移
TOC_FRAME_X, TOC_FRAME_CX, TOC_FRAME_CY = 6936105, 4287520, 522533
TOC_FRAME_Y_OFF = 10160


def toc_entries_from_source(prs: Presentation) -> list[str]:
    """Generic: pull section-like lines from the source TOC/outline slide."""
    # Prefer a slide that looks like a TOC (many x.y headings)
    best: list[str] = []
    for idx in range(min(len(prs.slides), 8)):
        lines = extract_texts(prs, idx)
        sections = []
        for line in lines:
            line = re.sub(r"\s+", " ", line).strip()
            # 2.1 / 3.2 / 10.1 等节号，或 「一、」一级条（较少作目录）
            if re.match(r"^\d+\.\d+", line):
                sections.append(line)
        if len(sections) > len(best):
            best = sections

    def key(s: str):
        m = re.match(r"(\d+)\.(\d+)", s)
        return (int(m.group(1)), int(m.group(2))) if m else (999, 999)

    seen = set()
    ordered = []
    for s in sorted(best, key=key):
        if s not in seen:
            seen.add(s)
            ordered.append(s)
    return ordered


def dedupe_toc_layout_logos(layouts_dir: Path) -> None:
    """Layout2/3/4 模板里同一 logo 被放了两次；只保留最下方一个。"""
    for name in ("slideLayout2.xml", "slideLayout3.xml", "slideLayout4.xml"):
        path = layouts_dir / name
        if not path.exists():
            continue
        xml = path.read_text(encoding="utf-8")
        rels_path = layouts_dir / "_rels" / f"{name}.rels"
        if not rels_path.exists():
            continue
        rels = rels_path.read_text(encoding="utf-8")
        # logo media typically image5.png
        logo_rids = set(
            re.findall(
                r'Id="(rId\d+)"[^>]*Target="../media/image5\.png"',
                rels,
            )
        )
        if not logo_rids:
            logo_rids = set(
                re.findall(
                    r'Target="../media/image5\.png"[^>]*Id="(rId\d+)"',
                    rels,
                )
            )
        if not logo_rids:
            continue

        pics = list(re.finditer(r"<p:pic>.*?</p:pic>", xml, re.S))
        logo_pics = []
        for m in pics:
            emb = re.search(r'r:embed="(rId\d+)"', m.group(0))
            if not emb or emb.group(1) not in logo_rids:
                continue
            ys = [int(y) for _, y in re.findall(r'<a:off x="(-?\d+)" y="(-?\d+)"/>', m.group(0))]
            logo_pics.append((max(ys) if ys else 0, m.start(), m.end()))
        if len(logo_pics) <= 1:
            continue
        # keep bottom-most logo (max y); remove others
        logo_pics.sort(key=lambda t: t[0], reverse=True)
        keep = logo_pics[0]
        to_remove = sorted(logo_pics[1:], key=lambda t: t[1], reverse=True)
        for _y, a, b in to_remove:
            xml = xml[:a] + xml[b:]
            print(f"  removed duplicate logo from {name} (kept y={keep[0]})")
        path.write_text(xml, encoding="utf-8")


def _toc_num_shape(shape_id: int, num: int, title_y: int) -> str:
    """克隆 Layout4 序号框：透明底 + #2576B5 描边 + #0070C0 数字。"""
    s = f"{num:02d}"
    y = title_y + TOC_NUM_Y_OFF
    return f'''<p:sp>
  <p:nvSpPr><p:cNvPr id="{shape_id}" name="目录序号{num}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{TOC_NUM_X}" y="{y}"/><a:ext cx="{TOC_NUM_CX}" cy="{TOC_NUM_CY}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln w="28575"><a:solidFill><a:srgbClr val="2576B5"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" anchor="ctr"/><a:lstStyle/>
    <a:p><a:pPr algn="ctr"/>
      <a:r><a:rPr lang="zh-CN" altLang="en-US" sz="2800" b="1">
        <a:solidFill><a:srgbClr val="0070C0"/></a:solidFill>
        <a:latin typeface="微软雅黑"/><a:ea typeface="微软雅黑"/>
      </a:rPr><a:t>{xml_text(s[0])}</a:t></a:r>
      <a:r><a:rPr lang="en-US" altLang="zh-CN" sz="2800" b="1">
        <a:solidFill><a:srgbClr val="0070C0"/></a:solidFill>
        <a:latin typeface="微软雅黑"/><a:ea typeface="微软雅黑"/>
      </a:rPr><a:t>{xml_text(s[1])}</a:t></a:r>
    </a:p>
  </p:txBody>
</p:sp>'''


def _toc_frame_shape(shape_id: int, title_y: int) -> str:
    """克隆 Layout4 标题描边框（空心，文字由旁侧文本框承担）。"""
    y = title_y + TOC_FRAME_Y_OFF
    return f'''<p:sp>
  <p:nvSpPr><p:cNvPr id="{shape_id}" name="目录标题框"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{TOC_FRAME_X}" y="{y}"/><a:ext cx="{TOC_FRAME_CX}" cy="{TOC_FRAME_CY}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
    <a:ln w="28575"><a:solidFill><a:srgbClr val="2576B5"/></a:solidFill></a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="zh-CN"/></a:p></p:txBody>
</p:sp>'''


def _toc_title_shape(shape_id: int, text: str, title_y: int) -> str:
    """克隆 Layout4 标题占位符外观：蓝字 #2576B5、微软雅黑加粗。"""
    return f'''<p:sp>
  <p:nvSpPr><p:cNvPr id="{shape_id}" name="目录标题"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{TOC_TITLE_X}" y="{title_y}"/><a:ext cx="{TOC_TITLE_CX}" cy="{TOC_TITLE_CY}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/><a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr vert="horz" lIns="90000" tIns="46800" rIns="90000" bIns="46800" anchor="ctr">
      <a:normAutofit/>
    </a:bodyPr>
    <a:lstStyle/>
    <a:p><a:pPr algn="l"/>
      <a:r><a:rPr lang="zh-CN" altLang="en-US" sz="2400" b="1">
        <a:solidFill><a:srgbClr val="2576B5"/></a:solidFill>
        <a:latin typeface="微软雅黑"/><a:ea typeface="微软雅黑"/>
      </a:rPr><a:t>{xml_text(text)}</a:t></a:r>
    </a:p>
  </p:txBody>
</p:sp>'''


def build_single_toc_slide(items: list[str]) -> tuple[str, str]:
    """Return (layout_file, slide_xml). 单页；前 5 条用版式，第 6 条起同构追加。"""
    n = len(items)
    if n <= 0:
        items = ["\u00a0"]
        n = 1

    # n 为 3/4/5：纯占位符（版式已去重 logo）
    if n in TOC_PH_IDX:
        idxs, layout = TOC_PH_IDX[n]
        parts = []
        sid = 2
        for i, idx in enumerate(idxs):
            parts.append(
                ph_sp(
                    sid,
                    f"目录项{idx}",
                    f'sz="half" idx="{idx}"',
                    para(items[i], bold=True, sz=2400, tight=True),
                )
            )
            sid += 1
        return layout, wrap_slide(parts)

    if n < 3:
        idxs, layout = TOC_PH_IDX[3]
        parts = []
        sid = 2
        for i, idx in enumerate(idxs):
            text = items[i] if i < n else "\u00a0"
            parts.append(
                ph_sp(
                    sid,
                    f"目录项{idx}",
                    f'sz="half" idx="{idx}"',
                    para(text, bold=True, sz=2400, tight=True),
                )
            )
            sid += 1
        return layout, wrap_slide(parts)

    # n > 5：前 5 条走 Layout4 占位符；第 6 条起只追加同构行（不改前五条）
    layout = "slideLayout4.xml"
    idxs = TOC_PH_IDX[5][0]
    parts = []
    sid = 2
    for i, idx in enumerate(idxs):
        parts.append(
            ph_sp(
                sid,
                f"目录项{idx}",
                f'sz="half" idx="{idx}"',
                para(items[i], bold=True, sz=2400, tight=True),
            )
        )
        sid += 1

    y = TOC_ROW_YS[-1] + TOC_ROW_GAP
    logo_top = 6283948
    for i, text in enumerate(items[5:]):
        num = i + 6
        if y + TOC_TITLE_CY > logo_top - 80000:
            break
        parts.append(_toc_num_shape(sid, num, y))
        sid += 1
        parts.append(_toc_frame_shape(sid, y))
        sid += 1
        parts.append(_toc_title_shape(sid, text, y))
        sid += 1
        y += TOC_ROW_GAP
    return layout, wrap_slide(parts)


def choose_layout(title: str, body: list[str], placements: list) -> str:
    # Keep Layout6 chrome (blue title bar) even on figure pages
    if re.match(r"^\d+\.\d", title) and not body and not placements:
        return "slideLayout5.xml"
    heads = [b for b in body if re.match(r"^[一二三四五]、", b)]
    if len(heads) >= 2 and not placements:
        return "slideLayout7.xml"
    return "slideLayout6.xml"


def build_content_slide(
    layout: str, title: str, body: list[str], image_items: list
) -> str:
    parts = []
    sid = 2
    has_img = bool(image_items)
    body_sz = SZ_BODY_IMG if has_img else SZ_BODY

    if layout == "slideLayout5.xml":
        m = re.match(r"^(2\.\d+)\s*(.*)$", title)
        chap = m.group(1) if m else ""
        t = (m.group(2).strip() if m and m.group(2).strip() else title)
        parts.append(
            ph_sp(sid, "标题", 'type="title" hasCustomPrompt="1"', para(t, bold=True, sz=SZ_SECTION))
        )
        sid += 1
        lines = ([chap] if chap else []) + body
        parts.append(
            ph_sp(
                sid,
                "正文",
                'type="body" idx="1" hasCustomPrompt="1"',
                "".join(para(x, bold=True, sz=2000, tight=True) for x in lines) or empty_para(),
            )
        )
        sid += 1
    elif layout == "slideLayout7.xml":
        split_at = len(body) // 2
        for i, b in enumerate(body):
            if i > 0 and re.match(r"^[二三四]、", b):
                split_at = i
                break
        left, right = body[:split_at] or body, body[split_at:]
        parts.append(ph_sp(sid, "标题", 'type="title"', para(title, bold=True, sz=SZ_TITLE)))
        sid += 1
        parts.append(
            ph_sp(
                sid,
                "左",
                'sz="half" idx="1"',
                "".join(para(x, sz=body_sz, tight=True) for x in left) or empty_para(),
            )
        )
        sid += 1
        parts.append(
            ph_sp(
                sid,
                "右",
                'sz="half" idx="2"',
                "".join(para(x, sz=body_sz, tight=True) for x in right) or empty_para(),
            )
        )
        sid += 1
    else:
        parts.append(ph_sp(sid, "标题", 'type="title"', para(title, bold=True, sz=SZ_TITLE)))
        sid += 1
        # Full body; tighter spacing so text pages are less empty
        parts.append(
            ph_sp(
                sid,
                "内容",
                'sz="half" idx="1"',
                "".join(para(x, sz=body_sz, tight=True) for x in body) or empty_para(),
            )
        )
        sid += 1

    for rid, fname, x, y, cx, cy in image_items:
        parts.append(make_pic(sid, fname, rid, x, y, cx, cy))
        sid += 1
    return wrap_slide(parts)


def write_slide(slides_dir: Path, rels_dir: Path, num: int, layout: str, xml: str, img_rels: list[tuple[str, str]]):
    name = f"slide{num}.xml"
    (slides_dir / name).write_text(xml, encoding="utf-8")
    lines = [
        f'<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/{layout}"/>'
    ]
    rid_n = 2
    mapping = []
    for dest, _geom in img_rels:
        rid = f"rId{rid_n}"
        lines.append(
            f'<Relationship Id="{rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/{dest}"/>'
        )
        mapping.append(rid)
        rid_n += 1
    (rels_dir / f"{name}.rels").write_text(
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">\n  '
        + "\n  ".join(lines)
        + "\n</Relationships>\n",
        encoding="utf-8",
    )
    return name, mapping


def convert_ppt_to_pptx(ppt_path: Path) -> Path:
    """Convert .ppt -> .pptx via PowerPoint COM."""
    out = ppt_path.with_name(ppt_path.stem + "_converted.pptx")
    if out.exists() and out.stat().st_mtime >= ppt_path.stat().st_mtime:
        return out
    import win32com.client  # type: ignore

    pp = win32com.client.Dispatch("PowerPoint.Application")
    try:
        try:
            pp.Visible = 1
        except Exception:
            pass
        pres = pp.Presentations.Open(str(ppt_path.resolve()), WithWindow=False)
        try:
            # 24 = ppSaveAsOpenXMLPresentation
            pres.SaveAs(str(out.resolve()), 24)
        finally:
            pres.Close()
    finally:
        try:
            pp.Quit()
        except Exception:
            pass
    return out


def convert_wmf_to_png(media_dir: Path, png_dir: Path) -> None:
    wmfs = sorted(media_dir.glob("*.wmf"))
    if not wmfs:
        return
    png_dir.mkdir(parents=True, exist_ok=True)
    import win32com.client  # type: ignore

    pp = win32com.client.Dispatch("PowerPoint.Application")
    try:
        try:
            pp.Visible = 1
        except Exception:
            pass
        for wmf in wmfs:
            png = png_dir / (wmf.stem + ".png")
            if png.exists() and png.stat().st_mtime >= wmf.stat().st_mtime:
                continue
            pres = pp.Presentations.Add()
            try:
                slide = pres.Slides.Add(1, 12)
                shape = slide.Shapes.AddPicture(str(wmf.resolve()), False, True, 10, 10)
                target_w = 15 * 72
                if shape.Width > 0:
                    ratio = target_w / float(shape.Width)
                    shape.Width = target_w
                    shape.Height = float(shape.Height) * ratio
                shape.Export(str(png.resolve()), 2)
                print("WMF->PNG", wmf.name, "->", png.name)
            except Exception as e:
                print("WMF fail", wmf.name, e)
            finally:
                try:
                    pres.Saved = True
                    pres.Close()
                except Exception:
                    pass
    finally:
        try:
            pp.Quit()
        except Exception:
            pass


def configure_paths(src: Path, out: Path | None = None) -> None:
    global SRC, OUT, WORK, PNG_DIR, SRC_UNPACKED
    src = src.resolve()
    if not src.exists():
        raise FileNotFoundError(src)
    if src.suffix.lower() == ".ppt":
        src = convert_ppt_to_pptx(src)
    elif src.suffix.lower() != ".pptx":
        raise ValueError(f"unsupported source: {src}")
    SRC = src
    stem = src.name.replace("_converted", "")
    if stem.lower().endswith(".pptx"):
        stem = Path(stem).stem
    OUT = out.resolve() if out else (ROOT / "output" / f"{stem}_培训课件_内容保真.pptx")
    WORK = ROOT / f"work_{stem}"
    PNG_DIR = ROOT / f"media_png_{stem}"
    SRC_UNPACKED = ROOT / f"src_unpacked_{stem}"



def main(src: Path | None = None, out: Path | None = None):
    if src is None:
        raise SystemExit("usage: python scripts/build_training_course.py <源.ppt|源.pptx>")
    configure_paths(Path(src), Path(out) if out else None)
    print("SRC", SRC)
    print("OUT", OUT)

    if SRC_UNPACKED.exists():
        shutil.rmtree(SRC_UNPACKED)
    SRC_UNPACKED.mkdir(parents=True)
    with zipfile.ZipFile(SRC) as z:
        z.extractall(SRC_UNPACKED)

    convert_wmf_to_png(SRC_UNPACKED / "ppt" / "media", PNG_DIR)

    if WORK.exists():
        shutil.rmtree(WORK)
    WORK.mkdir(parents=True)
    with zipfile.ZipFile(TPL) as z:
        z.extractall(WORK)

    # 目录版式自带双份 logo，先去重再填内容
    dedupe_toc_layout_logos(WORK / "ppt" / "slideLayouts")

    media_dir = WORK / "ppt" / "media"
    media_dir.mkdir(exist_ok=True)
    renamed = {}
    for f in (SRC_UNPACKED / "ppt" / "media").iterdir():
        dest = f"src_{f.name}"
        shutil.copy2(f, media_dir / dest)
        renamed[f.name] = dest
    # also stage png conversions
    if PNG_DIR.exists():
        for png in PNG_DIR.glob("*.png"):
            dest = f"src_{png.stem}.png"
            shutil.copy2(png, media_dir / dest)

    prs = Presentation(str(SRC))
    slides_dir = WORK / "ppt" / "slides"
    rels_dir = slides_dir / "_rels"
    rels_dir.mkdir(exist_ok=True)

    slide_files: list[str] = []

    def commit_slide(layout: str, xml: str, dest_geom: list[tuple[str, dict]]):
        name = f"slide{len(slide_files) + 1}.xml"
        rel_lines = [
            f'<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/{layout}"/>'
        ]
        for i, (dest, _g) in enumerate(dest_geom):
            rel_lines.append(
                f'<Relationship Id="rId{i+2}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/{dest}"/>'
            )
        (slides_dir / name).write_text(xml, encoding="utf-8")
        (rels_dir / f"{name}.rels").write_text(
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">\n  '
            + "\n  ".join(rel_lines)
            + "\n</Relationships>\n",
            encoding="utf-8",
        )
        slide_files.append(name)
        print(f"{name}: {layout} imgs={len(dest_geom)}")

    # --- Cover ---
    cover_lines = extract_texts(prs, 0, keep_org=True)
    cover_title = next((t for t in cover_lines if "导航原理" in t), None)
    if not cover_title:
        cover_title = cover_lines[0] if cover_lines else SRC.stem
    cover_body = [t for t in cover_lines if t != cover_title]
    commit_slide(
        "slideLayout1.xml",
        wrap_slide(
            [
                ph_sp(2, "标题", 'type="ctrTitle"', para(cover_title, bold=True, sz=SZ_COVER_TITLE)),
                ph_sp(
                    3,
                    "副标题",
                    'type="subTitle" idx="1"',
                    "".join(para(x, bold=True, sz=SZ_COVER_SUB) for x in cover_body)
                    or empty_para(),
                ),
            ]
        ),
        [],
    )

    # --- TOC: 单页；>5 条时同页追加与模板一致的「序号块+标题条」---
    toc_items = toc_entries_from_source(prs)
    toc_layout, toc_xml = build_single_toc_slide(toc_items)
    commit_slide(toc_layout, toc_xml, [])
    print("TOC layout:", toc_layout, "items:", toc_items)

    # --- Content = source slides 3..38 ---
    for src_i in range(2, len(prs.slides)):
        lines = extract_texts(prs, src_i)
        placements = fit_placements(extract_image_placements(src_i, SRC_UNPACKED))
        if not lines:
            title, body = "（本页图示）", []
        else:
            title = lines[0]
            for line in lines:
                if re.match(r"^[一二三四五六七八]、", line):
                    title = line
                    break
            body = [x for x in lines if x != title]

        layout = choose_layout(title, body, placements)
        dest_geom: list[tuple[str, dict]] = []
        for pl in placements:
            dest = resolve_media_file(pl["media"], media_dir, renamed)
            if dest:
                dest_geom.append((dest, pl))

        image_items = [
            (f"rId{i+2}", dest, g["x"], g["y"], g["cx"], g["cy"])
            for i, (dest, g) in enumerate(dest_geom)
        ]
        xml = build_content_slide(layout, title, body, image_items)
        commit_slide(layout, xml, dest_geom)
        print(f"  title={title[:28]!r}")

    # cleanup old template slides not in list
    for old in slides_dir.glob("slide*.xml"):
        if old.name not in slide_files:
            old.unlink(missing_ok=True)
            (rels_dir / f"{old.name}.rels").unlink(missing_ok=True)

    # presentation rels + sldIdLst
    pres_rels_path = WORK / "ppt" / "_rels" / "presentation.xml.rels"
    pres_rels = pres_rels_path.read_text(encoding="utf-8")
    rel_entries = re.findall(r"<Relationship\s[^>]*/>", pres_rels)
    kept = []
    for e in rel_entries:
        typ = re.search(r'Type="([^"]+)"', e)
        tgt = re.search(r'Target="([^"]+)"', e)
        if not typ or not tgt:
            continue
        if typ.group(1).endswith("/slide") or tgt.group(1).startswith("slides/"):
            continue
        kept.append(e)
    slide_rids = []
    next_id = 20
    for sn in slide_files:
        rid = f"rId{next_id}"
        next_id += 1
        kept.append(
            f'<Relationship Id="{rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/{sn}"/>'
        )
        slide_rids.append(rid)
    pres_rels_path.write_text(
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">\n  '
        + "\n  ".join(kept)
        + "\n</Relationships>\n",
        encoding="utf-8",
    )
    pres_path = WORK / "ppt" / "presentation.xml"
    pres_xml = pres_path.read_text(encoding="utf-8")
    sld_ids = "".join(
        f'<p:sldId id="{256+i}" r:id="{rid}"/>' for i, rid in enumerate(slide_rids)
    )
    pres_xml = re.sub(
        r"<p:sldIdLst>.*?</p:sldIdLst>",
        f"<p:sldIdLst>{sld_ids}</p:sldIdLst>",
        pres_xml,
        flags=re.S,
    )
    pres_path.write_text(pres_xml, encoding="utf-8")

    # content types
    ct_path = WORK / "[Content_Types].xml"
    ct = ct_path.read_text(encoding="utf-8")
    ct = re.sub(r'<Override PartName="/ppt/slides/slide\d+\.xml"[^/]*/>\s*', "", ct)
    if 'Extension="wmf"' not in ct:
        ct = ct.replace("</Types>", '  <Default Extension="wmf" ContentType="image/x-wmf"/>\n</Types>')
    overrides = "\n".join(
        f'  <Override PartName="/ppt/slides/{sn}" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        for sn in slide_files
    )
    ct = ct.replace("</Types>", overrides + "\n</Types>")
    ct_path.write_text(ct, encoding="utf-8")

    subprocess.check_call([sys.executable, str(PPTX_SCRIPTS / "clean.py"), str(WORK)])
    OUT.parent.mkdir(exist_ok=True)
    subprocess.check_call(
        [
            sys.executable,
            str(PPTX_SCRIPTS / "office" / "pack.py"),
            str(WORK),
            str(OUT),
            "--original",
            str(TPL),
        ]
    )
    print("Wrote", OUT, OUT.stat().st_size)


if __name__ == "__main__":
    import argparse

    ap = argparse.ArgumentParser(description="套用培训课件模板（内容保真）")
    ap.add_argument("src", help="源课件 .ppt 或 .pptx")
    ap.add_argument("-o", "--out", default=None, help="输出 pptx 路径（默认 output/{主名}_培训课件_内容保真.pptx）")
    args = ap.parse_args()
    main(Path(args.src), Path(args.out) if args.out else None)
